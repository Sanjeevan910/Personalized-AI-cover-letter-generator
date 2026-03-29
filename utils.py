import streamlit as st
from openai import OpenAI
import PyPDF2
from docx import Document
import io
import base64
import json
import csv
import openpyxl
from pptx import Presentation
from bs4 import BeautifulSoup
from config import API_KEY, BASE_URL, MODEL_NAME, JOB_FILE_TYPES, IMAGE_TYPES

try:
    import extract_msg
    HAS_MSG = True
except ImportError:
    HAS_MSG = False

client = OpenAI(
    api_key=API_KEY,
    base_url=BASE_URL
)

def extract_text_from_file(uploaded_file):
    text = ""
    if uploaded_file is None: return text
    try:
        if uploaded_file.name.endswith('.pdf'):
            # OCR is being used here for extracting text from PDF
            import pytesseract
            from pdf2image import convert_from_bytes
            images = convert_from_bytes(uploaded_file.read())
            for img in images: text += pytesseract.image_to_string(img) + "\n"
        elif uploaded_file.name.endswith(('.png', '.jpg', '.jpeg', '.webp')):
            # OCR is being used here for extracting text from images
            import pytesseract
            from PIL import Image
            img = Image.open(uploaded_file)
            text = pytesseract.image_to_string(img)
        elif uploaded_file.name.endswith('.docx'):
            doc = Document(uploaded_file)
            for para in doc.paragraphs: text += para.text + "\n"
        elif uploaded_file.name.endswith('.txt'):
            text = uploaded_file.read().decode("utf-8")
    except Exception as e:
        st.error(f"Error reading file: {e}")
        return None
    return text

def generate_cover_letter_api(company, role, resume_text, job_description, max_words):
    try:
        prompt = f"""
        Write a highly personalized cover letter for the position of {role} at {company}.
        
        INSTRUCTIONS:
        1. Analyze the Candidate's Resume below to extract their Name, Contact Details, and Key Skills. Use these details to personalize the letter header and content.
        2. Strictly adhere to a word limit of approximately {max_words} words.
        3. Tone: Formal, confident, and tailored. 
        4. Format: Standard business letter format (Header with Candidate Name/Contact -> Date -> Hiring Manager -> Body -> Sign-off).
        
        Candidate's Resume Context:
        {resume_text}
        
        Job Description:
        {job_description}
        """
        response = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[
                {"role": "system", "content": "You are an expert executive career coach."},
                {"role": "user", "content": prompt}
            ]
        )
        return response.choices[0].message.content
    except Exception as e:
        st.error(f"Error: {e}")
        return None

def _extract_text_from_job_file(uploaded_file):
    """Extract raw text from a job-posting file using OCR where appropriate."""
    name = uploaded_file.name.lower()
    ext = name.rsplit('.', 1)[-1] if '.' in name else ''

    if ext in IMAGE_TYPES:
        # OCR is being used here for extracting text from images
        import pytesseract
        from PIL import Image
        img = Image.open(uploaded_file)
        return pytesseract.image_to_string(img)

    try:
        if ext == 'pdf':
            # OCR is being used here for extracting text from PDF
            import pytesseract
            from pdf2image import convert_from_bytes
            images = convert_from_bytes(uploaded_file.read())
            return "\n".join(pytesseract.image_to_string(img) for img in images)

        elif ext in ('docx', 'doc'):
            doc = Document(uploaded_file)
            return "\n".join(p.text for p in doc.paragraphs)

        elif ext == 'txt':
            return uploaded_file.read().decode('utf-8', errors='replace')

        elif ext == 'html':
            soup = BeautifulSoup(uploaded_file.read(), 'html.parser')
            return soup.get_text(separator='\n')

        elif ext == 'json':
            data = json.load(uploaded_file)
            return json.dumps(data, indent=2)

        elif ext == 'csv':
            content = uploaded_file.read().decode('utf-8', errors='replace')
            reader = csv.reader(io.StringIO(content))
            return "\n".join(", ".join(row) for row in reader)

        elif ext == 'xlsx':
            wb = openpyxl.load_workbook(uploaded_file, data_only=True)
            lines = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    lines.append("  ".join(str(c) if c is not None else '' for c in row))
            return "\n".join(lines)

        elif ext == 'pptx':
            prs = Presentation(uploaded_file)
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        lines.append(shape.text)
            return "\n".join(lines)

        elif ext == 'eml':
            import email
            msg = email.message_from_bytes(uploaded_file.read())
            parts = []
            for part in msg.walk():
                if part.get_content_type() == 'text/plain':
                    parts.append(part.get_payload(decode=True).decode('utf-8', errors='replace'))
            return "\n".join(parts) if parts else msg.get_payload()

        elif ext == 'msg':
            if not HAS_MSG:
                st.error("`extract-msg` library is not installed. Run: pip install extract-msg")
                return ''
            raw = uploaded_file.read()
            msg_obj = extract_msg.Message(io.BytesIO(raw))
            return f"Subject: {msg_obj.subject}\n\n{msg_obj.body}"

        else:
            return uploaded_file.read().decode('utf-8', errors='replace')

    except Exception as e:
        st.error(f"Error reading file: {e}")
        return ''


def extract_job_details(job_file):
    """Universal extractor: supports text-based formats and images via OCR."""
    try:
        name = job_file.name.lower()
        ext = name.rsplit('.', 1)[-1] if '.' in name else ''

        json_prompt = """
Extract the following details from this job posting and return ONLY valid JSON:
{
    "company": "Company Name",
    "role": "Job Title / Role",
    "description": "Full Job Description / Requirements / Competencies"
}
If any field is missing, use \"Unknown\" or infer from context."""

        # Text extraction path using OCR
        raw_text = _extract_text_from_job_file(job_file)
        if not raw_text:
            st.error("Could not extract any text from the file.")
            return None
        full_prompt = f"{json_prompt}\n\n---\nFILE CONTENT:\n{raw_text[:12000]}"
        response = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[{"role": "user", "content": full_prompt}]
        )

        content = response.choices[0].message.content
        if "```json" in content:
            content = content.split("```json", 1)[1].split("```", 1)[0]
        elif "```" in content:
            content = content.split("```", 1)[1].split("```", 1)[0]
        return json.loads(content.strip())

    except Exception as e:
        st.error(f"Error extracting job details: {e}")
        return None
